#!/usr/bin/env python3
"""Build an executive assessment deck from findings.json (python-pptx).

Tracks ADO Story AB#5048.
"""
import sys
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

data = json.load(open(sys.argv[1]))
out = sys.argv[2]
prs = Presentation()
NAVY = RGBColor(0x1F, 0x4E, 0x78)


def title_slide(t, s):
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = t
    sl.placeholders[1].text = s


def bullets(title, lines):
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = title
    tf = sl.placeholders[1].text_frame
    tf.clear()
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = l
        p.font.size = Pt(16)


title_slide("Azure Landing Zone Assessment",
            "Current state vs. CAF / WAF — generated " + data["GeneratedOn"][:10])

bullets("Framework Scores",
        [f'{f["Framework"]}: {f["Score"]}/100' for f in data["Frameworks"]])

bullets("Area Scores",
        [f'{a["Framework"]} — {a["Area"]}: {a["Score"]}' for a in data["Areas"] if a["Score"] is not None])

def sev(g):
    # Defensive: a null/missing severity must not crash the whole deck (AB#5089).
    return (g.get("Severity") or "unknown").upper()

top = data["Gaps"][:10]
bullets("Top 10 Prioritized Gaps",
        [f'[{sev(g)}] {g["Area"]}: {g["Title"]}' for g in top])

bullets("Manual Review Items (questionnaire remainder)",
        [f'{m["Area"]}: {m["Title"]}' for m in data["Manual"][:12]])

prs.save(out)
